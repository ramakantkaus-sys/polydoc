"""PPTX writer, built on python-pptx.

Slides are a fixed canvas, so the interesting problem is the inverse of the reader's:
turning a flowing document into slides that fit.

When the document already has :class:`~polydoc.model.Slide` blocks (it came from a
presentation) they are written directly. Otherwise the body is *paginated*: each
top-level heading starts a slide, and content is split further when it would overflow
the estimated text height. Without that, converting a long report produces one slide
with text running off the bottom.
"""

from __future__ import annotations

from typing import Any, BinaryIO, Dict, List, Optional, Sequence, Tuple

from ...model import (
    Alignment,
    Block,
    BlockContainer,
    CodeBlock,
    Container,
    Document,
    Heading,
    HorizontalRule,
    Image,
    Inline,
    InlineImage,
    LineBreak,
    Link,
    ListBlock,
    ListItem,
    ListStyle,
    Page,
    PageBreak,
    Paragraph,
    Quote,
    Section,
    Slide,
    Table,
    TableCell,
    Text,
    TextStyle,
)
from ..base import Writer, require
from ..registry import register_writer

__all__ = ["PptxWriter"]

_EMU_PER_POINT = 12700

_ALIGNMENT_NAMES = {
    Alignment.LEFT: "LEFT",
    Alignment.CENTER: "CENTER",
    Alignment.RIGHT: "RIGHT",
    Alignment.JUSTIFY: "JUSTIFY",
}

#: Layout indices in the default python-pptx template.
_LAYOUT_TITLE = 0
_LAYOUT_TITLE_CONTENT = 1
_LAYOUT_TITLE_ONLY = 5
_LAYOUT_BLANK = 6

#: Rough vertical cost of one line at the default body size, in points.
_LINE_HEIGHT = 26.0
#: Characters that fit on one line of the content placeholder, approximately.
_CHARS_PER_LINE = 58

#: Universal list styles mapped onto PowerPoint auto-numbering schemes.
_AUTONUM_TYPES = {
    ListStyle.ORDERED: "arabicPeriod",
    ListStyle.LOWER_ALPHA: "alphaLcPeriod",
    ListStyle.UPPER_ALPHA: "alphaUcPeriod",
    ListStyle.LOWER_ROMAN: "romanLcPeriod",
    ListStyle.UPPER_ROMAN: "romanUcPeriod",
}

#: Bullet glyphs per nesting depth, matching PowerPoint's own defaults.
_BULLET_CHARS = ("\u2022", "\u25e6", "\u25aa")


@register_writer
class PptxWriter(Writer):
    """Writes PowerPoint presentations (``.pptx``)."""

    format = "pptx"
    extensions = (".pptx",)
    aliases = ("powerpoint", "slides")
    mime_types = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    extra = "pptx"
    description = "Microsoft PowerPoint (Open XML), with automatic pagination"

    def write(self, document: Document, stream: BinaryIO, **options: Any) -> None:
        pptx = require("pptx", "Writing PPTX", extra="pptx", package="python-pptx")
        builder = _PptxBuilder(pptx, document, options)
        builder.build().save(stream)


class _PptxBuilder:
    def __init__(self, pptx_module: Any, document: Document, options: Dict[str, Any]) -> None:
        self.pptx = pptx_module
        self.document = document
        self.options = options
        self.embed_images: bool = options.get("embed_images", True)
        self.max_lines: int = options.get("max_lines_per_slide", 11)
        self.include_notes: bool = options.get("include_notes", True)

        template = options.get("template")
        self.native = pptx_module.Presentation(template) if template else pptx_module.Presentation()
        self._apply_geometry()

    # -- setup ----------------------------------------------------------------
    def _apply_geometry(self) -> None:
        from pptx.util import Pt

        geometry = self.document.geometry
        widescreen = self.options.get("widescreen")
        if geometry is not None and self.options.get("template") is None:
            self.native.slide_width = Pt(geometry.size.width)
            self.native.slide_height = Pt(geometry.size.height)
        elif widescreen and self.options.get("template") is None:
            self.native.slide_width = Pt(960)
            self.native.slide_height = Pt(540)

    @property
    def _content_width(self) -> float:
        return self.native.slide_width / _EMU_PER_POINT

    @property
    def _content_height(self) -> float:
        return self.native.slide_height / _EMU_PER_POINT

    # -- entry point ----------------------------------------------------------
    def build(self) -> Any:
        self._apply_metadata()
        slides = self.document.slides
        if slides:
            for slide in slides:
                self._write_slide(slide.title, slide.content, slide.notes)
        else:
            for title, blocks in self._paginate():
                self._write_slide(title, blocks, None)

        if not self.native.slides:
            # Never emit a zero-slide deck; PowerPoint treats it as damaged.
            self._write_slide(self.document.metadata.title or "Untitled", [], None)
        return self.native

    def _apply_metadata(self) -> None:
        """Copy our metadata across, clearing python-pptx's template defaults.

        As with DOCX, the stock template carries an author and creation date that are
        not ours; inheriting them would misattribute the file.
        """
        import datetime as _datetime

        meta = self.document.metadata
        props = self.native.core_properties

        props.title = meta.title or ""
        props.author = "; ".join(meta.authors) if meta.authors else ""
        props.subject = meta.subject or ""
        props.keywords = ", ".join(meta.keywords) if meta.keywords else ""
        props.comments = meta.description or ""
        props.last_modified_by = ""

        stamp = self.options.get("timestamp", True)
        now = _datetime.datetime.now(_datetime.timezone.utc).replace(tzinfo=None)
        if meta.created or stamp:
            props.created = meta.created or now
        if meta.modified or stamp:
            props.modified = meta.modified or now

    # -- pagination -----------------------------------------------------------
    def _paginate(self) -> List[Tuple[Optional[str], List[Block]]]:
        """Split a flowing document into (title, blocks) pairs, one per slide."""
        flattened = self._flatten(self.document.body)
        slides: List[Tuple[Optional[str], List[Block]]] = []

        title: Optional[str] = None
        current: List[Block] = []
        budget = self.max_lines

        def commit() -> None:
            nonlocal current, budget
            if title is not None or current:
                slides.append((title, current))
            current = []
            budget = self.max_lines

        for block in flattened:
            if isinstance(block, Heading) and block.level <= 2:
                commit()
                title = block.text
                continue
            if isinstance(block, PageBreak):
                commit()
                continue

            cost = self._cost(block)
            # A block taller than a whole slide has to go somewhere; give it its own.
            if current and budget - cost < 0:
                pending_title = title
                commit()
                title = pending_title
            current.append(block)
            budget -= cost

        commit()

        if not slides and self.document.metadata.title:
            slides.append((self.document.metadata.title, []))
        return slides

    def _flatten(self, blocks: Sequence[Block]) -> List[Block]:
        """Unwrap containers so pagination sees a flat stream."""
        out: List[Block] = []
        for block in blocks:
            if isinstance(block, Section):
                if block.title:
                    out.append(Heading(list(block.title), max(1, min(6, block.level or 1))))
                out.extend(self._flatten(block.content))
            elif isinstance(block, (Page, Container)) and not isinstance(block, ListItem):
                if isinstance(block, Container) and block.role == "sheet" and block.name:
                    out.append(Heading.of(block.name, 2))
                out.extend(self._flatten(block.content))
                if isinstance(block, Page):
                    out.append(PageBreak())
            else:
                out.append(block)
        return out

    def _cost(self, block: Block) -> int:
        """Estimated line count, used to decide when a slide is full."""
        if isinstance(block, ListBlock):
            return sum(
                max(1, -(-len(item.text) // _CHARS_PER_LINE)) + sum(
                    self._cost(sub) for sub in item.sublists
                )
                for item in block.items
            )
        if isinstance(block, Table):
            return len(block.rows) + 1
        if isinstance(block, CodeBlock):
            return len(block.lines) + 1
        if isinstance(block, Image):
            return self.max_lines  # an image gets a slide to itself
        if isinstance(block, HorizontalRule):
            return 1
        text = block.text
        return max(1, -(-len(text) // _CHARS_PER_LINE))

    # -- slide construction ---------------------------------------------------
    def _write_slide(
        self,
        title: Optional[str],
        blocks: Sequence[Block],
        notes: Optional[str],
    ) -> Any:
        body_blocks = [b for b in blocks if not isinstance(b, PageBreak)]
        tables = [b for b in body_blocks if isinstance(b, Table)]
        images = [b for b in body_blocks if isinstance(b, Image)]
        text_blocks = [b for b in body_blocks if not isinstance(b, (Table, Image))]

        layout_index = self._choose_layout(title, text_blocks, tables, images)
        slide = self.native.slides.add_slide(self.native.slide_layouts[layout_index])

        if title is not None:
            self._set_title(slide, title)

        if text_blocks:
            placeholder = self._body_placeholder(slide)
            if placeholder is not None:
                self._fill_text_frame(placeholder.text_frame, text_blocks)
            else:
                self._add_textbox(slide, text_blocks)

        top = self._next_top(slide, bool(text_blocks))
        for table in tables:
            top = self._add_table(slide, table, top)
        for image in images:
            top = self._add_image(slide, image, top)

        if notes and self.include_notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except (AttributeError, KeyError):  # pragma: no cover
                pass
        return slide

    def _choose_layout(
        self,
        title: Optional[str],
        text_blocks: Sequence[Block],
        tables: Sequence[Block],
        images: Sequence[Block],
    ) -> int:
        if title is not None and not text_blocks and not tables and not images:
            return _LAYOUT_TITLE
        if text_blocks:
            return _LAYOUT_TITLE_CONTENT
        if title is not None:
            return _LAYOUT_TITLE_ONLY
        return _LAYOUT_BLANK

    def _set_title(self, slide: Any, title: str) -> None:
        placeholder = None
        try:
            placeholder = slide.shapes.title
        except (AttributeError, KeyError):
            placeholder = None
        if placeholder is None:
            from pptx.util import Pt

            box = slide.shapes.add_textbox(
                Pt(40), Pt(24), Pt(self._content_width - 80), Pt(56)
            )
            frame = box.text_frame
            frame.text = title
            frame.paragraphs[0].runs[0].font.size = Pt(30)
            frame.paragraphs[0].runs[0].font.bold = True
            return
        placeholder.text_frame.text = title

    def _body_placeholder(self, slide: Any) -> Any:
        """Find the content placeholder, skipping any title placeholder.

        Compared by placeholder *type* rather than object identity: python-pptx wraps
        the underlying XML in a new Python object on each access, so ``shape is
        slide.shapes.title`` is never true and an identity check would hand back the
        title placeholder itself.
        """
        from pptx.enum.shapes import PP_PLACEHOLDER

        title_types = {
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
            PP_PLACEHOLDER.VERTICAL_TITLE,
        }
        for shape in slide.placeholders:
            try:
                if shape.placeholder_format.type in title_types:
                    continue
            except (AttributeError, KeyError, ValueError):
                pass
            if shape.has_text_frame:
                return shape
        return None

    def _add_textbox(self, slide: Any, blocks: Sequence[Block]) -> Any:
        from pptx.util import Pt

        box = slide.shapes.add_textbox(
            Pt(40), Pt(110), Pt(self._content_width - 80), Pt(self._content_height - 150)
        )
        box.text_frame.word_wrap = True
        self._fill_text_frame(box.text_frame, blocks)
        return box

    def _next_top(self, slide: Any, has_text: bool) -> float:
        """Vertical cursor for shapes added after the placeholders."""
        if not has_text:
            return 120.0
        lowest = 120.0
        for shape in slide.shapes:
            top = getattr(shape, "top", None)
            height = getattr(shape, "height", None)
            if top is None or height is None:
                continue
            lowest = max(lowest, (top + height) / _EMU_PER_POINT)
        return min(lowest + 12.0, self._content_height - 80)

    # -- text -----------------------------------------------------------------
    def _fill_text_frame(self, frame: Any, blocks: Sequence[Block]) -> None:
        frame.word_wrap = True
        frame.clear()
        first = True
        for block in blocks:
            first = self._write_block(frame, block, first=first, level=0)

    def _write_block(self, frame: Any, block: Block, first: bool, level: int) -> bool:
        from pptx.util import Pt

        if isinstance(block, Paragraph):
            paragraph = self._paragraph(frame, first)
            self._apply_paragraph(paragraph, block, level)
            self._runs(paragraph, block.content)
            return False

        if isinstance(block, Heading):
            paragraph = self._paragraph(frame, first)
            paragraph.level = level
            self._runs(paragraph, block.content, extra=TextStyle(bold=True))
            return False

        if isinstance(block, ListBlock):
            return self._write_list(frame, block, first, level)

        if isinstance(block, CodeBlock):
            current = first
            for line in block.code.split("\n"):
                paragraph = self._paragraph(frame, current)
                paragraph.level = min(8, level)
                run = paragraph.add_run()
                run.text = line or " "
                run.font.name = "Consolas"
                run.font.size = Pt(14)
                current = False
            return False

        if isinstance(block, Quote):
            current = first
            for child in block.content:
                current = self._write_block(frame, child, current, level + 1)
            if block.attribution:
                paragraph = self._paragraph(frame, False)
                paragraph.level = min(8, level + 1)
                run = paragraph.add_run()
                run.text = f"\u2014 {block.attribution}"
                run.font.italic = True
            return False

        if isinstance(block, HorizontalRule):
            paragraph = self._paragraph(frame, first)
            paragraph.level = level
            run = paragraph.add_run()
            run.text = "\u2500" * 24
            return False

        if isinstance(block, BlockContainer):
            current = first
            for child in block.content:
                current = self._write_block(frame, child, current, level)
            return current

        if block.text:
            paragraph = self._paragraph(frame, first)
            paragraph.level = level
            run = paragraph.add_run()
            run.text = block.text
            return False
        return first

    def _write_list(self, frame: Any, block: ListBlock, first: bool, level: int) -> bool:
        current = first
        for index, item in enumerate(block.items):
            # PowerPoint has no checkbox, so task state is shown with a glyph.
            prefix = ""
            if item.checked is not None:
                prefix = "\u2612 " if item.checked else "\u2610 "

            wrote_prefix = False
            for child in item.content:
                if isinstance(child, ListBlock):
                    current = self._write_list(frame, child, current, level + 1)
                    continue
                if isinstance(child, Paragraph):
                    paragraph = self._paragraph(frame, current)
                    paragraph.level = min(8, level)
                    content: List[Inline] = list(child.content)
                    if prefix and not wrote_prefix:
                        content = [Text(prefix)] + content
                        wrote_prefix = True
                    self._runs(paragraph, content)
                    # Real bullet properties, not a text prefix, so the marker style
                    # survives being read back.
                    self._apply_bullet(
                        paragraph,
                        block.marker_style,
                        block.start if index == 0 else None,
                        level,
                    )
                    current = False
                else:
                    current = self._write_block(frame, child, current, level + 1)
        return current

    def _apply_bullet(
        self,
        paragraph: Any,
        marker_style: ListStyle,
        start: Optional[int],
        level: int,
    ) -> None:
        """Write ``a:buChar`` / ``a:buAutoNum`` onto a paragraph's properties.

        python-pptx exposes no bullet API, so the DrawingML is built directly. Element
        order inside ``a:pPr`` is schema-significant, hence the insertion before
        ``a:defRPr`` rather than a plain append.
        """
        from pptx.oxml.ns import qn
        from lxml import etree

        p_pr = paragraph._p.get_or_add_pPr()

        for tag in ("a:buNone", "a:buAutoNum", "a:buChar", "a:buFont"):
            for existing in p_pr.findall(qn(tag)):
                p_pr.remove(existing)

        indent_unit = 285750  # 22.5pt in EMU, PowerPoint's default step
        p_pr.set("marL", str(indent_unit * (level + 1)))
        p_pr.set("indent", str(-indent_unit))

        elements: List[Any] = []
        if marker_style is ListStyle.NONE:
            elements.append(etree.SubElement(p_pr, qn("a:buNone")))
        elif marker_style.is_ordered:
            font = etree.Element(qn("a:buFont"))
            font.set("typeface", "+mj-lt")
            elements.append(font)
            auto = etree.Element(qn("a:buAutoNum"))
            auto.set("type", _AUTONUM_TYPES.get(marker_style, "arabicPeriod"))
            if start is not None and start > 1:
                auto.set("startAt", str(start))
            elements.append(auto)
        else:
            font = etree.Element(qn("a:buFont"))
            font.set("typeface", "Arial")
            font.set("panose", "020B0604020202020204")
            font.set("pitchFamily", "34")
            font.set("charset", "0")
            elements.append(font)
            char = etree.Element(qn("a:buChar"))
            char.set("char", _BULLET_CHARS[level % len(_BULLET_CHARS)])
            elements.append(char)

        # buNone was created in place by SubElement; the rest need positioning.
        anchor = p_pr.find(qn("a:defRPr"))
        for element in elements:
            if element.getparent() is p_pr:
                continue
            if anchor is not None:
                anchor.addprevious(element)
            else:
                p_pr.append(element)

    @staticmethod
    def _paragraph(frame: Any, first: bool) -> Any:
        """Reuse the frame's initial empty paragraph, then append."""
        if first and frame.paragraphs:
            return frame.paragraphs[0]
        return frame.add_paragraph()

    def _apply_paragraph(self, paragraph: Any, block: Paragraph, level: int) -> None:
        from pptx.enum.text import PP_ALIGN

        paragraph.level = min(8, level)
        alignment = block.style.alignment
        if alignment is not None:
            name = _ALIGNMENT_NAMES.get(alignment)
            if name is not None:
                try:
                    paragraph.alignment = getattr(PP_ALIGN, name)
                except AttributeError:  # pragma: no cover
                    pass

    def _runs(
        self,
        paragraph: Any,
        content: Sequence[Inline],
        extra: Optional[TextStyle] = None,
    ) -> None:
        for node in content:
            if isinstance(node, Text):
                self._run(paragraph, node.text, node.style.merge(extra) if extra else node.style)
            elif isinstance(node, Link):
                for child in node.content:
                    style = getattr(child, "style", TextStyle())
                    run = self._run(paragraph, child.text, style)
                    if run is not None and node.href:
                        try:
                            run.hyperlink.address = node.href
                        except (AttributeError, ValueError):  # pragma: no cover
                            pass
            elif isinstance(node, LineBreak):
                try:
                    paragraph.add_line_break()
                except AttributeError:  # pragma: no cover
                    self._run(paragraph, "\n", TextStyle())
            elif isinstance(node, InlineImage):
                if node.alt:
                    self._run(paragraph, f"[{node.alt}]", TextStyle(italic=True))
            elif node.text:
                self._run(paragraph, node.text, TextStyle())

    def _run(self, paragraph: Any, text: str, style: TextStyle) -> Any:
        from pptx.dml.color import RGBColor
        from pptx.util import Pt

        if not text:
            return None
        run = paragraph.add_run()
        run.text = text
        font = run.font
        if style.bold is not None:
            font.bold = style.bold
        if style.italic is not None:
            font.italic = style.italic
        if style.underline is not None:
            font.underline = style.underline
        if style.is_monospace:
            font.name = style.font_family or "Consolas"
        elif style.font_family:
            font.name = style.font_family
        if style.font_size:
            font.size = Pt(style.font_size)
        if style.color:
            try:
                value = style.color.lstrip("#")
                font.color.rgb = RGBColor(
                    int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)
                )
            except (ValueError, IndexError):
                pass
        return run

    # -- tables and images ----------------------------------------------------
    def _add_table(self, slide: Any, block: Table, top: float) -> float:
        from pptx.util import Pt

        rows, columns = block.dimensions
        if not rows or not columns:
            return top

        margin = 40.0
        width = self._content_width - 2 * margin
        height = min(rows * 26.0, self._content_height - top - 30)
        shape = slide.shapes.add_table(
            rows, columns, Pt(margin), Pt(top), Pt(width), Pt(max(height, 26.0))
        )
        native = shape.table
        native.first_row = bool(block.header_rows)

        if block.column_widths:
            total = sum(w for w in block.column_widths if w > 0)
            if total > 0:
                for index, column_width in enumerate(block.column_widths[:columns]):
                    if column_width > 0:
                        native.columns[index].width = Pt(width * column_width / total)

        for row_index, row in enumerate(block.rows[:rows]):
            column = 0
            for cell in row.cells:
                if column >= columns:
                    break
                target = native.cell(row_index, column)
                span_end = min(columns - 1, column + cell.colspan - 1)
                row_end = min(rows - 1, row_index + cell.rowspan - 1)
                if span_end > column or row_end > row_index:
                    try:
                        target.merge(native.cell(row_end, span_end))
                    except (ValueError, KeyError):  # pragma: no cover
                        pass
                self._fill_cell(target, cell)
                column = span_end + 1

        return top + (shape.height / _EMU_PER_POINT) + 14.0

    def _fill_cell(self, native_cell: Any, cell: TableCell) -> None:
        from pptx.util import Pt

        frame = native_cell.text_frame
        frame.clear()
        blocks = [b for b in cell.content if not isinstance(b, (Table, Image))]
        if not blocks:
            frame.text = ""
            return
        first = True
        for block in blocks:
            first = self._write_block(frame, block, first=first, level=0)
        for paragraph in frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size is None:
                    run.font.size = Pt(14)

    def _add_image(self, slide: Any, block: Image, top: float) -> float:
        from io import BytesIO

        from pptx.util import Pt

        if not (self.embed_images and block.data):
            if block.alt or block.caption:
                self._add_textbox(slide, [Paragraph.of(f"[Image: {block.alt or block.caption}]")])
            return top

        margin = 40.0
        available_width = self._content_width - 2 * margin
        available_height = self._content_height - top - 30

        width = block.width or available_width
        height = block.height or available_height
        # Preserve aspect ratio inside the available box.
        scale = min(available_width / width, available_height / height, 1.0)
        width *= scale
        height *= scale
        left = (self._content_width - width) / 2

        try:
            picture = slide.shapes.add_picture(
                BytesIO(block.data), Pt(left), Pt(top), Pt(width), Pt(height)
            )
        except Exception:
            return top

        bottom = top + (picture.height / _EMU_PER_POINT)
        if block.caption:
            box = slide.shapes.add_textbox(
                Pt(margin), Pt(bottom + 6), Pt(available_width), Pt(24)
            )
            frame = box.text_frame
            frame.text = block.caption
            for run in frame.paragraphs[0].runs:
                run.font.size = Pt(12)
                run.font.italic = True
            bottom += 30
        return bottom + 12.0
