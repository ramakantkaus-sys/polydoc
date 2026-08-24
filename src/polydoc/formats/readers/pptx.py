"""PPTX reader, built on python-pptx.

A slide is a bag of positioned shapes, not a document flow, so two decisions matter:

**Ordering.** Shapes are stored in z-order, which is unrelated to reading order. They
are sorted top-to-bottom then left-to-right using their frame positions, so a
two-column slide reads sensibly.

**Roles.** The title placeholder becomes :attr:`Slide.title`; other text frames become
paragraphs and lists. Body text uses ``paragraph.level`` for nesting, which is genuine
structure rather than something we have to infer.

Speaker notes land on :attr:`Slide.notes`, and grouped shapes are flattened recursively.
"""

from __future__ import annotations

from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple

from ...model import (
    Alignment,
    BBox,
    Block,
    Container,
    Document,
    Heading,
    Image,
    Inline,
    LineBreak,
    Link,
    ListBlock,
    ListItem,
    ListStyle,
    PageGeometry,
    Paragraph,
    ParagraphStyle,
    Size,
    Slide,
    Table,
    TableCell,
    TableRow,
    Text,
    TextStyle,
    merge_runs,
)
from ..base import Reader, require
from ..registry import register_reader
from ..source import Source

__all__ = ["PptxReader"]

#: EMU per point. python-pptx measures everything in English Metric Units.
_EMU_PER_POINT = 12700.0


@register_reader
class PptxReader(Reader):
    """Reads PowerPoint presentations (``.pptx``, ``.pptm``)."""

    format = "pptx"
    extensions = (".pptx", ".pptm", ".ppsx")
    aliases = ("powerpoint", "slides")
    mime_types = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    extra = "pptx"
    archive_based = True
    description = "Microsoft PowerPoint (Open XML)"

    def read(self, source: Source, **options: Any) -> Document:
        self.enforce_limits(source, **options)
        pptx = require("pptx", "Reading PPTX", extra="pptx", package="python-pptx")
        handle = source.open_path()
        native = pptx.Presentation(str(handle) if hasattr(handle, "suffix") else handle)

        parser = _PptxParser(
            pptx,
            extract_images=options.get("extract_images", True),
            include_notes=options.get("include_notes", True),
            heading_from_title=options.get("heading_from_title", False),
        )

        document = Document()
        geometry = self._geometry(native)
        document.geometry = geometry

        for index, native_slide in enumerate(native.slides, start=1):
            document.append(parser.slide(native_slide, index, geometry))

        self._read_metadata(native, document)
        if document.metadata.title is None and document.slides:
            document.metadata.title = document.slides[0].heading

        return self.finalise(document, source)

    @staticmethod
    def _geometry(native: Any) -> Optional[PageGeometry]:
        width = getattr(native, "slide_width", None)
        height = getattr(native, "slide_height", None)
        if not width or not height:
            return None
        return PageGeometry(
            size=Size(width / _EMU_PER_POINT, height / _EMU_PER_POINT),
            margin_left=0.0,
            margin_right=0.0,
            margin_top=0.0,
            margin_bottom=0.0,
        )

    @staticmethod
    def _read_metadata(native: Any, document: Document) -> None:
        props = native.core_properties
        meta = document.metadata
        meta.title = props.title or None
        if props.author:
            meta.authors = [p.strip() for p in props.author.split(";") if p.strip()]
        meta.subject = props.subject or None
        if props.keywords:
            import re

            meta.keywords = [p.strip() for p in re.split(r"[;,]", props.keywords) if p.strip()]
        meta.description = props.comments or None
        meta.category = props.category or None
        meta.language = props.language or None
        meta.created = props.created
        meta.modified = props.modified


class _PptxParser:
    def __init__(
        self,
        pptx_module: Any,
        extract_images: bool = True,
        include_notes: bool = True,
        heading_from_title: bool = False,
    ) -> None:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        self.pptx = pptx_module
        self.MSO = MSO_SHAPE_TYPE
        self.extract_images = extract_images
        self.include_notes = include_notes
        self.heading_from_title = heading_from_title

    # -- slides ---------------------------------------------------------------
    def slide(self, native: Any, index: int, geometry: Optional[PageGeometry]) -> Slide:
        title: Optional[str] = None
        blocks: List[Block] = []

        shapes = self._ordered_shapes(native.shapes)
        for shape, is_title in shapes:
            if is_title and title is None:
                text = self._shape_text(shape)
                if text:
                    title = text
                    if self.heading_from_title:
                        blocks.append(Heading.of(text, 2))
                    continue
            blocks.extend(self.shape(shape))

        layout = None
        try:
            layout = native.slide_layout.name
        except AttributeError:  # pragma: no cover
            layout = None

        notes = self._notes(native) if self.include_notes else None

        return Slide(
            content=blocks,
            title=title,
            layout=layout,
            notes=notes,
            index=index,
            geometry=geometry,
        )

    def _ordered_shapes(self, shapes: Any) -> List[Tuple[Any, bool]]:
        """Sort shapes into reading order, flagging the title placeholder.

        Shapes are stored in z-order; a slide's visual reading order is roughly
        top-to-bottom then left-to-right, so we sort on the frame origin. Shapes
        without a position (rare) keep their original relative order at the end.
        """
        positioned: List[Tuple[float, float, int, Any, bool]] = []
        unpositioned: List[Tuple[Any, bool]] = []

        for order, shape in enumerate(shapes):
            is_title = self._is_title(shape)
            top = getattr(shape, "top", None)
            left = getattr(shape, "left", None)
            if top is None or left is None:
                unpositioned.append((shape, is_title))
                continue
            # Quantise the vertical position so shapes roughly side by side sort
            # left to right rather than by a few EMU of jitter.
            band = round(top / (_EMU_PER_POINT * 12))
            positioned.append((band, left, order, shape, is_title))

        positioned.sort(key=lambda entry: (entry[0], entry[1], entry[2]))
        ordered = [(shape, is_title) for _, _, _, shape, is_title in positioned]
        ordered.extend(unpositioned)

        # The title always reads first regardless of where it sits on the slide.
        titles = [entry for entry in ordered if entry[1]]
        others = [entry for entry in ordered if not entry[1]]
        return titles + others

    def _is_title(self, shape: Any) -> bool:
        try:
            if not shape.is_placeholder:
                return False
            from pptx.enum.shapes import PP_PLACEHOLDER

            kind = shape.placeholder_format.type
            return kind in (
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
                PP_PLACEHOLDER.VERTICAL_TITLE,
            )
        except (AttributeError, KeyError, ValueError):
            return False

    def _notes(self, native: Any) -> Optional[str]:
        try:
            if not native.has_notes_slide:
                return None
            frame = native.notes_slide.notes_text_frame
        except (AttributeError, KeyError):
            return None
        if frame is None:
            return None
        text = frame.text.strip()
        return text or None

    @staticmethod
    def _shape_text(shape: Any) -> Optional[str]:
        try:
            if not shape.has_text_frame:
                return None
        except AttributeError:
            return None
        text = shape.text_frame.text.strip()
        return " ".join(text.split()) if text else None

    # -- shapes ---------------------------------------------------------------
    def shape(self, shape: Any) -> List[Block]:
        shape_type = getattr(shape, "shape_type", None)

        if shape_type == self.MSO.GROUP:
            out: List[Block] = []
            for child, _ in self._ordered_shapes(shape.shapes):
                out.extend(self.shape(child))
            if not out:
                return []
            return [Container(out, role="group")]

        if shape_type in (self.MSO.PICTURE, self.MSO.LINKED_PICTURE):
            image = self._picture(shape)
            return [image] if image is not None else []

        try:
            if shape.has_table:
                return [self._table(shape)]
        except AttributeError:
            pass

        try:
            if shape.has_chart:
                return self._chart(shape)
        except (AttributeError, ValueError):
            pass

        try:
            if shape.has_text_frame:
                return self._text_frame(shape)
        except AttributeError:
            pass

        return []

    def _bbox(self, shape: Any) -> Optional[BBox]:
        top = getattr(shape, "top", None)
        left = getattr(shape, "left", None)
        width = getattr(shape, "width", None)
        height = getattr(shape, "height", None)
        if None in (top, left, width, height):
            return None
        x0 = left / _EMU_PER_POINT
        y0 = top / _EMU_PER_POINT
        return BBox(x0, y0, x0 + width / _EMU_PER_POINT, y0 + height / _EMU_PER_POINT)

    def _picture(self, shape: Any) -> Optional[Image]:
        data = None
        mime = None
        src = None
        try:
            image = shape.image
            mime = image.content_type
            src = image.filename or f"image{image.ext and '.' + image.ext or ''}"
            if self.extract_images:
                data = image.blob
        except (AttributeError, ValueError, KeyError):
            # A linked picture has no embedded blob.
            pass

        alt = ""
        try:
            alt = shape._element._nvXxPr.cNvPr.get("descr") or shape.name or ""
        except AttributeError:
            alt = getattr(shape, "name", "") or ""

        box = self._bbox(shape)
        return Image(
            src=src,
            data=data,
            alt=alt,
            width=box.width if box else None,
            height=box.height if box else None,
            mime_type=mime,
            bbox=box,
        )

    def _chart(self, shape: Any) -> List[Block]:
        """Represent a chart by its underlying data, which is what survives export."""
        try:
            chart = shape.chart
        except (AttributeError, ValueError):
            return []

        title = None
        try:
            if chart.has_title:
                title = chart.chart_title.text_frame.text.strip() or None
        except (AttributeError, ValueError):
            title = None

        try:
            categories = [str(c) for c in chart.plots[0].categories]
        except (AttributeError, IndexError, ValueError):
            categories = []

        rows: List[List[Any]] = []
        try:
            series = list(chart.series)
        except (AttributeError, ValueError):
            series = []

        if categories and series:
            rows.append(["Category"] + [s.name or f"Series {i + 1}" for i, s in enumerate(series)])
            for index, category in enumerate(categories):
                row: List[Any] = [category]
                for s in series:
                    values = list(s.values)
                    row.append(values[index] if index < len(values) else "")
                rows.append(row)

        if not rows:
            return [Paragraph.of(f"[Chart: {title or shape.name}]")]

        table = Table.from_rows(rows, header=True, caption=title or "Chart data")
        table.attrs["role"] = "chart"
        return [table]

    def _text_frame(self, shape: Any) -> List[Block]:
        frame = shape.text_frame
        box = self._bbox(shape)
        entries: List[Tuple[int, ListStyle, List[Inline], ParagraphStyle]] = []

        for paragraph in frame.paragraphs:
            content = self._paragraph_content(paragraph)
            if not "".join(node.text for node in content).strip():
                continue
            level = int(getattr(paragraph, "level", 0) or 0)
            style = self._paragraph_style(paragraph)
            bullet = self._bullet_style(paragraph)
            entries.append((level, bullet, content, style))

        if not entries:
            return []

        blocks = self._assemble(entries)
        if box is not None:
            for block in blocks:
                if block.bbox is None:
                    block.bbox = box
        return blocks

    def _assemble(
        self,
        entries: Sequence[Tuple[int, ListStyle, List[Inline], ParagraphStyle]],
    ) -> List[Block]:
        """Group indented paragraphs into nested lists, leaving level-0 prose alone."""
        out: List[Block] = []
        pending: List[Tuple[int, ListStyle, List[Inline], ParagraphStyle]] = []

        def flush() -> None:
            if not pending:
                return
            out.extend(_build_list(pending))
            pending.clear()

        for level, bullet, content, style in entries:
            is_list = bullet is not None and (level > 0 or bullet is not ListStyle.NONE)
            if is_list:
                pending.append((level, bullet, content, style))
            else:
                flush()
                out.append(Paragraph(content, style))
        flush()
        return out

    def _bullet_style(self, paragraph: Any) -> Optional[ListStyle]:
        """Read the bullet definition from the paragraph's XML properties.

        python-pptx exposes no bullet API, so we look for ``a:buChar`` (a glyph),
        ``a:buAutoNum`` (numbering) or ``a:buNone`` directly.
        """
        try:
            p_pr = paragraph._pPr
        except AttributeError:
            return None
        if p_pr is None:
            # No explicit properties: an indented paragraph inherits a bullet.
            return ListStyle.BULLET if int(getattr(paragraph, "level", 0) or 0) > 0 else ListStyle.NONE

        namespace = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        if p_pr.find(f"{namespace}buNone") is not None:
            return ListStyle.NONE
        auto = p_pr.find(f"{namespace}buAutoNum")
        if auto is not None:
            return _AUTONUM_STYLES.get(auto.get("type", ""), ListStyle.ORDERED)
        if p_pr.find(f"{namespace}buChar") is not None:
            return ListStyle.BULLET
        return ListStyle.BULLET if int(getattr(paragraph, "level", 0) or 0) > 0 else ListStyle.NONE

    def _paragraph_style(self, paragraph: Any) -> ParagraphStyle:
        alignment = None
        try:
            if paragraph.alignment is not None:
                alignment = Alignment.coerce(paragraph.alignment.name)
        except AttributeError:
            alignment = None
        return ParagraphStyle(alignment=alignment)

    def _paragraph_content(self, paragraph: Any) -> List[Inline]:
        out: List[Inline] = []
        for run in paragraph.runs:
            style = self._run_style(run)
            href = None
            try:
                href = run.hyperlink.address
            except (AttributeError, KeyError):
                href = None
            text = run.text
            if not text:
                continue
            # PowerPoint uses a vertical tab for a soft line break inside a run.
            parts = text.split("\v")
            runs: List[Inline] = []
            for index, part in enumerate(parts):
                if index:
                    runs.append(LineBreak())
                if part:
                    runs.append(Text(part, style))
            if href:
                out.append(Link(runs, href=href))
            else:
                out.extend(runs)
        return merge_runs(out)

    def _run_style(self, run: Any) -> TextStyle:
        font = run.font
        colour = None
        try:
            if font.color is not None and font.color.type is not None:
                rgb = getattr(font.color, "rgb", None)
                if rgb is not None:
                    colour = f"#{str(rgb).lower()}"
        except (AttributeError, ValueError, KeyError):
            colour = None
        return TextStyle(
            bold=font.bold,
            italic=font.italic,
            underline=True if font.underline else None,
            font_family=font.name,
            font_size=font.size.pt if font.size is not None else None,
            color=colour,
        )

    # -- tables ---------------------------------------------------------------
    def _table(self, shape: Any) -> Table:
        native = shape.table
        rows: List[TableRow] = []
        first_row_header = bool(getattr(native, "first_row", False))

        # Track spans so continuation cells are skipped rather than duplicated.
        for row_index, native_row in enumerate(native.rows):
            cells: List[TableCell] = []
            for column_index, native_cell in enumerate(native_row.cells):
                if getattr(native_cell, "is_spanned", False):
                    continue
                content: List[Block] = []
                try:
                    content = self._text_frame(native_cell)
                except AttributeError:
                    content = []
                if not content and native_cell.text.strip():
                    content = [Paragraph.of(native_cell.text.strip())]
                cells.append(
                    TableCell(
                        content,
                        colspan=int(getattr(native_cell, "span_width", 1) or 1),
                        rowspan=int(getattr(native_cell, "span_height", 1) or 1),
                    )
                )
            if cells:
                rows.append(TableRow(cells, is_header=first_row_header and row_index == 0))

        table = Table(rows)
        box = self._bbox(shape)
        if box is not None:
            table.bbox = box
        widths = []
        try:
            widths = [round(c.width / _EMU_PER_POINT, 2) for c in native.columns]
        except (AttributeError, TypeError):
            widths = []
        if widths and any(widths):
            table.column_widths = widths
        return table


#: PowerPoint auto-numbering schemes mapped onto universal list styles.
_AUTONUM_STYLES = {
    "arabicPeriod": ListStyle.ORDERED,
    "arabicParenR": ListStyle.ORDERED,
    "arabicParenBoth": ListStyle.ORDERED,
    "alphaLcPeriod": ListStyle.LOWER_ALPHA,
    "alphaLcParenR": ListStyle.LOWER_ALPHA,
    "alphaUcPeriod": ListStyle.UPPER_ALPHA,
    "alphaUcParenR": ListStyle.UPPER_ALPHA,
    "romanLcPeriod": ListStyle.LOWER_ROMAN,
    "romanUcPeriod": ListStyle.UPPER_ROMAN,
}


def _build_list(
    entries: Sequence[Tuple[int, ListStyle, List[Inline], ParagraphStyle]],
) -> List[Block]:
    """Assemble ``(level, style, content, paragraph_style)`` tuples into nested lists."""
    if not entries:
        return []

    base_level = entries[0][0]
    root = ListBlock(marker_style=entries[0][1] or ListStyle.BULLET)
    stack: List[Tuple[int, ListBlock]] = [(base_level, root)]

    for level, bullet, content, style in entries:
        while len(stack) > 1 and level < stack[-1][0]:
            stack.pop()
        current_level, current = stack[-1]

        if level > current_level and current.items:
            sub = ListBlock(marker_style=bullet or ListStyle.BULLET, level=len(stack))
            parent_item = current.items[-1]
            parent_item.content.append(sub)
            parent_item.adopt(sub)
            stack.append((level, sub))
            current = sub

        if not current.items and bullet is not None and bullet is not ListStyle.NONE:
            current.marker_style = bullet

        item = ListItem([Paragraph(content, style)])
        current.items.append(item)
        current.adopt(item)

    return [root]
